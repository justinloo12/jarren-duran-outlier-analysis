#!/usr/bin/env python3
"""Red Sox outfield-logjam strategy deck (python-pptx, no external tooling).

Team-level reframe of the Duran analysis: the OF is a surplus asset; the
pitching is a strength; spend the surplus on bats. Embeds pre-rendered charts
from ../figures/.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
OUT = ROOT / "outputs" / "Red_Sox_Outfield_Strategy.pptx"

# palette
NAVY = "0C2340"; RED = "BD3039"; INK = "20232A"; MUTED = "6B7280"
BG = "F4F5F7"; WHITE = "FFFFFF"; GREEN = "2E7D32"; GOLD = "B7892F"; LINE = "D7DAE0"
HEAD = "Cambria"; BODY = "Calibri"

W, H = 13.333, 7.5
ASPECT = {"01": 1.60, "04": 1.60, "07": 10/5.8, "08": 10/6.4,
          "09": 1504/649, "10": 2087/679}


def C(hexs):
    return RGBColor.from_string(hexs)


def txt(slide, x, y, w, h, runs, size=16, color=INK, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
        space_after=4, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for t, o in para:
            r = p.add_run(); r.text = t
            f = r.font
            f.size = Pt(o.get("size", size)); f.bold = o.get("bold", bold)
            f.italic = o.get("italic", italic); f.name = o.get("font", font)
            f.color.rgb = C(o.get("color", color))
    return tb


def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if line:
        sh.line.color.rgb = C(line); sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def dot(slide, x, y, d, fill=RED, text=None, tcolor=WHITE, size=15):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill); sh.line.fill.background()
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True; r.font.name = HEAD
        r.font.color.rgb = C(tcolor)
    return sh


def image(slide, key, x, y, max_w, max_h, name):
    a = ASPECT[key]
    w = max_w; hh = w / a
    if hh > max_h:
        hh = max_h; w = hh * a
    ix = x + (max_w - w) / 2
    slide.shapes.add_picture(str(FIG / name), In(ix), In(y), In(w), In(hh))
    return ix, y, w, hh


def eyebrow(slide, text, x=0.62, y=0.42, color=RED):
    txt(slide, x, y, 11, 0.3, [[(text.upper(), {"color": color, "bold": True,
        "size": 12.5, "font": BODY})]], space_after=0)


def title(slide, text, x=0.62, y=0.72, w=12.1, size=30, color=NAVY):
    txt(slide, x, y, w, 0.95, [[(text, {"size": size, "bold": True,
        "font": HEAD, "color": color})]], space_after=0, line_spacing=1.0)


def footer(slide, dark=False):
    c = "9AA3B2" if dark else MUTED
    txt(slide, 0.62, 7.12, 12.1, 0.3,
        [[("Boston Red Sox · Outfield Strategy", {"color": c, "size": 9}),
          ("      Analysis: Statcast / FanGraphs / MLB Stats API · salary Spotrac · July 2026",
           {"color": c, "size": 9})]], space_after=0)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(color)


prs = Presentation()
prs.slide_width = In(W); prs.slide_height = In(H)
blank = prs.slide_layouts[6]


def new(color=BG):
    s = prs.slides.add_slide(blank); bg(s, color); return s


# ---------------------------------------------------------------- 1 TITLE
s = new(NAVY)
dot(s, 0.62, 0.62, 0.34, fill=RED)
txt(s, 1.05, 0.6, 10, 0.4, [[("BOSTON RED SOX  ·  ROSTER STRATEGY",
    {"color": "C9D2E3", "bold": True, "size": 13})]], space_after=0)
txt(s, 0.62, 2.15, 12.1, 2.0, [
    [("Turning the Outfield Logjam", {"size": 46, "bold": True, "font": HEAD,
      "color": WHITE})],
    [("into an Asset", {"size": 46, "bold": True, "font": HEAD, "color": RED})],
], space_after=2, line_spacing=1.0)
txt(s, 0.62, 4.35, 11.6, 1.0, [[(
    "Five bats for four spots. The pitching is carrying the team — so the fix "
    "isn't to buy arms, it's to convert a redundant outfield into the offense "
    "and youth the roster actually needs.",
    {"size": 17, "color": "D3D9E6"})]], line_spacing=1.12)
txt(s, 0.62, 6.6, 12, 0.4, [[("A data-driven roster case study · July 2026",
    {"color": "8C97AC", "size": 12, "italic": True})]], space_after=0)

# ---------------------------------------------------------------- 2 SITUATION
s = new()
eyebrow(s, "The setup")
title(s, "A down year for the bats — but a rare surplus in the grass")
stats = [("~70", "win pace in 2026 (37-48)", "the offense, not the pitching, is the drag"),
         ("Top-10", "starting rotation", "8th in WAR, 9th in ERA (3.76) — a strength"),
         ("5 → 4", "OF/DH for four spots", "the one area of genuine roster depth")]
cw, gap = 3.83, 0.31
for i, (big, lab, sub) in enumerate(stats):
    x = 0.62 + i * (cw + gap)
    card(s, x, 2.05, cw, 2.35)
    txt(s, x + 0.28, 2.28, cw - 0.5, 0.9, [[(big, {"size": 44, "bold": True,
        "font": HEAD, "color": RED if i == 0 else NAVY})]], space_after=0)
    txt(s, x + 0.3, 3.28, cw - 0.55, 0.4, [[(lab, {"size": 15, "bold": True,
        "color": NAVY})]], space_after=0)
    txt(s, x + 0.3, 3.66, cw - 0.55, 0.7, [[(sub, {"size": 12.5,
        "color": MUTED})]], line_spacing=1.05)
txt(s, 0.62, 4.75, 12.1, 1.4, [
    [("The Sox aren't buying their way to a 2026 title. ",
      {"bold": True, "color": INK, "size": 15}),
     ("They're setting up 2027 around a strong rotation and a cheap young "
      "outfield core. Right now the jam is even ",
      {"color": INK, "size": 15}),
     ("masked by Roman Anthony's injury", {"bold": True, "color": RED, "size": 15}),
     (" (60-day IL) — which makes it a 2027 problem to solve this winter, not a "
      "deadline scramble.", {"color": INK, "size": 15})],
], line_spacing=1.2)
footer(s)

# ---------------------------------------------------------------- 3 KEEP CORE
s = new()
eyebrow(s, "Keep")
title(s, "The young, cheap, controllable core — that's the 2027 outfield")
core = [("Roman Anthony", "LF · age 22", "3.8", "$2.6M · 8yr/$130M ext",
         "Franchise bat. Untouchable. On the IL now, back in 2027."),
        ("Ceddanne Rafaela", "CF · age 25", "3.2", "$2.0M · 8yr/$50M",
         "Elite center-field glove + SS/utility flex. Premium position, "
         "controlled cheap."),
        ("Wilyer Abreu", "RF · age 27", "2.8", "$0.8M · pre-arb → 2029",
         "Best surplus on the roster: plus RF glove, above-average LH bat.")]
cw, gap = 3.83, 0.31
for i, (nm, pa, war, sal, note) in enumerate(core):
    x = 0.62 + i * (cw + gap)
    card(s, x, 1.95, cw, 3.9)
    txt(s, x + 0.3, 2.2, cw - 0.55, 0.5, [[(nm, {"size": 19, "bold": True,
        "font": HEAD, "color": NAVY})]], space_after=0)
    txt(s, x + 0.3, 2.72, cw - 0.55, 0.3, [[(pa, {"size": 12.5, "bold": True,
        "color": RED})]], space_after=0)
    txt(s, x + 0.3, 3.25, cw - 0.55, 0.9, [[(war, {"size": 40, "bold": True,
        "font": HEAD, "color": GREEN})]], space_after=0)
    txt(s, x + 0.3, 4.18, cw - 0.55, 0.3, [[("projected WAR / yr",
        {"size": 11, "color": MUTED})]], space_after=0)
    txt(s, x + 0.3, 4.55, cw - 0.55, 0.35, [[(sal, {"size": 12, "bold": True,
        "color": INK})]], space_after=0)
    txt(s, x + 0.3, 4.95, cw - 0.55, 0.85, [[(note, {"size": 11.5,
        "color": MUTED})]], line_spacing=1.06)
txt(s, 0.62, 6.05, 12.1, 0.5, [[("Three spots, locked. Everything else flexes "
    "around them.", {"italic": True, "size": 14, "color": NAVY})]],
    space_after=0)
footer(s)

# ---------------------------------------------------------------- 4 MOVE
s = new()
eyebrow(s, "Move")
title(s, "The two that don't fit: a redundant bat and a dead contract")
cw = 5.86
data = [(0.62, "Jarren Duran", "LF/CF · 29 · $7.7M", RED,
         [("Movable piece.", "Controllable through 2028; above-average at true "
           "talent (~110 wRC+). Redundant behind the core and the priciest of "
           "the young group."),
          ("The asset to sell.", "The one everyday, controllable OF a "
           "contender will pay real prospects for.")]),
        (0.62 + cw + 0.33, "Masataka Yoshida", "DH · 32 · $18.6M", MUTED,
         [("Dead money.", "$18.6M through 2027 for a DH-only, below-average "
           "bat. Negative trade value."),
          ("Absorb, don't chase.", "Let it expire as a platoon DH, or move "
           "only in a partial salary-dump. Don't pay to escape it.")])]
for x, nm, meta, accent, rows in data:
    card(s, x, 1.95, cw, 3.95)
    txt(s, x + 0.32, 2.22, cw - 0.6, 0.5, [[(nm, {"size": 21, "bold": True,
        "font": HEAD, "color": NAVY})]], space_after=0)
    txt(s, x + 0.32, 2.78, cw - 0.6, 0.32, [[(meta, {"size": 13, "bold": True,
        "color": accent})]], space_after=0)
    yy = 3.35
    for head, body in rows:
        dot(s, x + 0.32, yy + 0.05, 0.13, fill=accent)
        txt(s, x + 0.58, yy, cw - 0.95, 1.2, [[
            (head + " ", {"bold": True, "color": INK, "size": 13.5}),
            (body, {"color": MUTED, "size": 13})]], line_spacing=1.08)
        yy += 1.25
footer(s)

# ---------------------------------------------------------------- 5 DURAN LUCK
s = new()
eyebrow(s, "Don't sell low")
title(s, "Duran's 2026 is unlucky — not a true collapse")
ix, iy, iw, ih = image(s, "01", 0.5, 1.75, 6.9, 4.55,
                       "01_babip_vs_league.png")
tx = 7.7
rows = [("2024 was earned, not a fluke", RED,
         "His BABIP (.344) sat close to his contact-quality xBABIP (.333); "
         "xwOBA .340 was legit. A good year, only modestly lucky."),
        ("2026 is the real anomaly — downward", NAVY,
         "wOBA .265 sits BELOW its xwOBA .286; BABIP .244 vs .311 expected. "
         "Real skill erosion, but the line is worse than the player."),
        ("True talent ≈ his 2025 (~110 wRC+)", GREEN,
         "An above-average regular. Both the 2024 high and 2026 low are tails "
         "around that.")]
yy = 1.95
for head, col, body in rows:
    dot(s, tx, yy + 0.04, 0.16, fill=col)
    txt(s, tx + 0.32, yy, 5.0, 1.3, [[(head, {"bold": True, "size": 14.5,
        "color": NAVY})], [(body, {"size": 12.8, "color": MUTED})]],
        line_spacing=1.08, space_after=2)
    yy += 1.38
txt(s, tx, yy + 0.05, 5.0, 0.6, [[("Selling now cashes the slump. Value hinges "
    "on the rebound.", {"italic": True, "bold": True, "size": 13,
    "color": RED})]], line_spacing=1.05)
footer(s)

# ---------------------------------------------------------------- 5b REBOUND SIM
s = new()
eyebrow(s, "The rebound math")
title(s, "10,000 simulations: the rebound is likely — unless the erosion is real")
image(s, "09", 0.4, 1.7, 8.3, 4.1, "09_rebound_probability.png")
tx = 9.0
txt(s, tx, 1.8, 3.7, 0.35, [[("P(REST-OF-SEASON wRC+ ≥ 100)",
    {"size": 11.5, "bold": True, "color": MUTED})]], space_after=0)
txt(s, tx, 2.15, 3.7, 0.85, [
    [("69%", {"size": 34, "bold": True, "font": HEAD, "color": GREEN}),
     ("  healthy prior", {"size": 12.5, "color": MUTED})],
    [("47%", {"size": 34, "bold": True, "font": HEAD, "color": RED}),
     ("  erosion-blended", {"size": 12.5, "color": MUTED})],
], space_after=2, line_spacing=1.0)
txt(s, tx, 3.85, 3.7, 0.35, [[("EXPECTED OFFSEASON TRADE VALUE",
    {"size": 11.5, "bold": True, "color": MUTED})]], space_after=0)
txt(s, tx, 4.2, 3.7, 0.8, [
    [("$22M", {"size": 24, "bold": True, "font": HEAD, "color": NAVY}),
     (" vs ", {"size": 13, "color": MUTED}),
     ("$18M", {"size": 24, "bold": True, "font": HEAD, "color": NAVY})],
    [("either way, well above the ~$10M a July sale fetches",
      {"size": 11.5, "color": MUTED})],
], space_after=2, line_spacing=1.05)
txt(s, 0.62, 6.15, 12.1, 0.7, [[("The honest line: ", {"size": 13.5,
    "bold": True, "color": RED}),
    ("the 22-point probability gap is the price of the worse chase / whiff / "
     "hard-hit rates. Hold-through-2026 survives the erosion stress test — "
     "but with less confidence in the rebound narrative itself.",
     {"size": 13.5, "color": INK})]], line_spacing=1.12)
footer(s)

# ------------------------------------------------------- 5c EROSION VERDICT
s = new()
eyebrow(s, "The tiebreaker")
title(s, "The bat speed says it's the approach — not the body")
image(s, "10", 0.5, 1.7, 12.33, 3.2, "10_erosion_decomposition.png")
cards = [("74.5 mph", "avg bat speed, 2026 — UP from 72.7 in '24", GREEN,
          "A declining hitter swings slower. Duran is swinging harder, "
          "longer and steeper — pressing, not fading."),
         ("29% whiff vs 95+", "up from 16% in '24 — but not bat-death", RED,
          "Whiffs on a rising bat speed = selling out. The chase leak is "
          "velo/offspeed; breaking-ball chase actually improved."),
         ("60 / 40", "re-weighted scenario odds → 60% rebound", NAVY,
          "Approach is fixable (see his own 2022→23 chase overhaul); the "
          "tools a buyer prices are intact.")]
cw, gap = 3.83, 0.31
for i, (big, lab, col, body) in enumerate(cards):
    x = 0.62 + i * (cw + gap)
    card(s, x, 5.05, cw, 1.85)
    txt(s, x + 0.26, 5.2, cw - 0.5, 0.5, [[(big, {"size": 21, "bold": True,
        "font": HEAD, "color": col}),
        ("  " + lab, {"size": 10.5, "color": MUTED})]], space_after=0,
        line_spacing=1.0)
    txt(s, x + 0.26, 5.78, cw - 0.5, 1.0, [[(body, {"size": 10.8,
        "color": INK})]], line_spacing=1.07)
footer(s)

# ---------------------------------------------------------------- 6 AGE CURVE
s = new()
eyebrow(s, "Context")
title(s, "Aging doesn't explain the drop — regression and luck do")
image(s, "04", 3.05, 1.65, 7.25, 4.5, "04_age_curve_overlay.png")
txt(s, 0.62, 6.35, 12.1, 0.7, [[("2024 → 2026 wRC+ fell ~70 points; a normal "
    "age 27→29 curve predicts only ~7. ", {"size": 14, "color": INK}),
    ("The gap is mean-reversion from an inflated peak plus a 2026 bad-luck "
     "tail — not skill cliff.", {"size": 14, "bold": True, "color": NAVY})]],
    line_spacing=1.1)
footer(s)

# ---------------------------------------------------------------- 7 WHAT TO GET
s = new()
eyebrow(s, "What to get back")
title(s, "Pitching is the strength — spend the surplus on offense")
txt(s, 0.62, 1.72, 12.1, 0.7, [[("With a top-10 rotation (3.76 ERA) already "
    "carrying the club, trading bats for arms is backwards. Convert the "
    "outfield surplus into what's thin:", {"size": 14.5, "color": INK})]],
    line_spacing=1.12)
needs = [("Controllable bats", GREEN,
          "Young, cost-controlled position players to fix an inconsistent "
          "lineup and lengthen it around the core."),
         ("Best talent / upside", NAVY,
          "In a retool you take the best player, not the positional patch — "
          "prospects and MLB-ready youth over rentals."),
         ("Salary relief", GOLD,
          "Attach Yoshida money where possible; free payroll for the 2027 "
          "window.")]
cw, gap = 3.83, 0.31
for i, (head, col, body) in enumerate(needs):
    x = 0.62 + i * (cw + gap)
    card(s, x, 2.7, cw, 2.75)
    dot(s, x + 0.3, 2.98, 0.42, fill=col, text=str(i + 1), size=16)
    txt(s, x + 0.3, 3.62, cw - 0.55, 0.45, [[(head, {"size": 16, "bold": True,
        "font": HEAD, "color": NAVY})]], space_after=0)
    txt(s, x + 0.3, 4.12, cw - 0.55, 1.2, [[(body, {"size": 12.8,
        "color": MUTED})]], line_spacing=1.1)
txt(s, 0.62, 5.75, 12.1, 0.5, [[("Keep the arms. Cash the outfield. ",
    {"bold": True, "italic": True, "size": 14, "color": NAVY}),
    ("Build the 2027 lineup, don't rent the 2026 one.",
     {"italic": True, "size": 14, "color": MUTED})]], space_after=0)
footer(s)

# ---------------------------------------------------------------- 8 MARKET
s = new()
eyebrow(s, "The market")
title(s, "Who pays: contenders with a hole in the outfield")
image(s, "08", 0.4, 1.6, 8.2, 4.7, "08_trade_fit_targets.png")
tx = 8.95
txt(s, tx, 1.85, 4.0, 0.5, [[("Best-fit buyers", {"size": 15, "bold": True,
    "font": HEAD, "color": NAVY})]], space_after=0)
for i, (t, d) in enumerate([("Phillies", "contender, 88 OF wRC+"),
                            ("Rays", ".602, weak OF"),
                            ("Marlins", "in the hunt, thin OF"),
                            ("Guardians", "contender, OF need")]):
    yy = 2.4 + i * 0.62
    dot(s, tx, yy + 0.03, 0.14, fill=RED)
    txt(s, tx + 0.28, yy, 3.9, 0.5, [[(t + "  ", {"bold": True, "size": 14,
        "color": NAVY}), (d, {"size": 12, "color": MUTED})]], space_after=0)
txt(s, tx, 5.15, 4.0, 1.1, [[("Trade a redundant bat to a team that needs one "
    "— and take back the ", {"size": 12.5, "color": INK}),
    ("hitting/youth", {"size": 12.5, "bold": True, "color": GREEN}),
    (" Boston is short on.", {"size": 12.5, "color": INK})]],
    line_spacing=1.12)
footer(s)

# ---------------------------------------------------------------- 9 PLAN
s = new()
eyebrow(s, "The sequence")
title(s, "Solve it this winter, from a position of strength")
steps = [("1", "Hold through 2026", "No contention urgency and Anthony hurt — "
          "keep Duran's bat, let his value rebound off an unlucky first half."),
         ("2", "Deal Duran in the offseason", "Anthony's healthy (jam returns), "
          "value restored. Trade him — plus Yoshida money — for controllable "
          "bats, not arms."),
         ("3", "Run it back in 2027", "OF = Anthony / Rafaela / Abreu; strong "
          "rotation intact; DH opens after 2027. A cheaper, younger, more "
          "balanced roster.")]
yy = 1.95
for num, head, body in steps:
    dot(s, 0.7, yy + 0.02, 0.5, fill=RED, text=num, size=18)
    txt(s, 1.4, yy, 6.7, 1.1, [[(head, {"size": 17, "bold": True, "font": HEAD,
        "color": NAVY})], [(body, {"size": 13.5, "color": MUTED})]],
        line_spacing=1.08, space_after=2)
    yy += 1.28
card(s, 8.55, 1.95, 4.15, 3.65, fill="0C2340", line=None)
txt(s, 8.85, 2.2, 3.6, 0.4, [[("THE RETURN", {"size": 12.5, "bold": True,
    "color": "C9D2E3"})]], space_after=0)
txt(s, 8.85, 2.72, 3.6, 1.0, [
    [("Sell now (nadir)", {"size": 13, "bold": True, "color": "F0B8BC"})],
    [("~$10–15M · a 45-FV bat", {"size": 14, "color": WHITE})]],
    space_after=2, line_spacing=1.05)
txt(s, 8.85, 3.85, 3.6, 1.1, [
    [("Sell after rebound", {"size": 13, "bold": True, "color": "A9E5B5"})],
    [("~$28M · a 50-FV, top-100", {"size": 14, "color": WHITE})],
    [("controllable talent", {"size": 14, "color": WHITE})]],
    space_after=2, line_spacing=1.05)
txt(s, 8.85, 5.05, 3.6, 0.5, [[("One grade of return = the reason to wait.",
    {"size": 11.5, "italic": True, "color": "C9D2E3"})]], line_spacing=1.05)
footer(s)

# ---------------------------------------------------------------- 10 CLOSER
s = new(NAVY)
dot(s, 0.62, 0.7, 0.34, fill=RED)
txt(s, 1.05, 0.68, 10, 0.4, [[("BOTTOM LINE", {"color": "C9D2E3",
    "bold": True, "size": 13})]], space_after=0)
lines = [("Keep the strength.", "A top-10 rotation and a cheap, young OF core "
          "are the foundation — don't trade into them."),
         ("Cash the surplus, not the slump.", "Move Duran this winter after a "
          "value rebound, for bats and youth — sell the 2025 player, not the "
          "unlucky 2026 one."),
         ("Quarantine the dead money.", "Absorb Yoshida; free the 2027 payroll "
          "and the DH slot.")]
yy = 1.95
for head, body in lines:
    dot(s, 0.7, yy + 0.06, 0.18, fill=RED)
    txt(s, 1.05, yy, 11.4, 1.0, [[(head + "  ", {"bold": True, "size": 21,
        "font": HEAD, "color": WHITE}), (body, {"size": 15, "color": "CDD4E2"})]],
        line_spacing=1.1, space_after=2)
    yy += 1.35
txt(s, 1.05, 6.15, 11.4, 0.8, [[("Confidence: ", {"bold": True, "size": 12.5,
    "color": "9AA6BC"}), ("high that 2024 shouldn't be the anchor and pitching "
    "is a strength; the open question is how much of 2026's skill dip (chase, "
    "whiff, hard-hit are all worse) persists — the one thing to watch before "
    "selling.", {"size": 12.5, "color": "9AA6BC", "italic": True})]],
    line_spacing=1.12)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
